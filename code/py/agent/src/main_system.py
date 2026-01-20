import os
import glob
import logging
import json
import os
import sys
import glob
import logging
import json

# Add project root to sys.path to allow 'from src...' imports when running directly
# This is necessary because we are using absolute imports (src.xxx) for IDE compatibility
# but running the script from inside the package structure.
current_dir = os.path.dirname(os.path.abspath(__file__))
project_root = os.path.dirname(current_dir)
if project_root not in sys.path:
    sys.path.append(project_root)

from src.extractors.information_extractor import InformationExtractor
from src.builders.graph_builder import GraphBuilder
from src.storages.vector_store import VectorStore
from src.reasoners.logic_reasoner import LogicReasoner

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class KnowledgeSystem:
    def __init__(self, storage_dir="data"):
        """
        Initialize the Knowledge System.
        """
        self.storage_dir = storage_dir
        self.extractor = InformationExtractor()
        self.graph_builder = GraphBuilder()
        self.vector_store = VectorStore()
        self.reasoner = LogicReasoner(self.graph_builder)
        self.processed_files = set() # Registry of processed file paths
        
        # Ensure storage directory exists
        os.makedirs(self.storage_dir, exist_ok=True)
        
        # Load existing data if available
        self.load_data()

    def load_documents(self, doc_dir):
        """
        Load documents from the specified directory.
        Supports .txt, .md, .csv, .json, .pdf, .xlsx, .pptx, .docx, .zip.
        Skips files that have already been processed.
        """
        files = glob.glob(os.path.join(doc_dir, "**/*.*"), recursive=True)
        logger.info(f"Found {len(files)} files in {doc_dir}")
        
        new_files_count = 0
        for file_path in files:
            # Normalize path
            abs_path = os.path.abspath(file_path)
            
            if abs_path in self.processed_files:
                logger.info(f"Skipping already processed file: {abs_path}")
                continue
                
            try:
                content = self._read_file(abs_path)
                if content:
                    self.process_content(content, abs_path)
                    self.processed_files.add(abs_path)
                    new_files_count += 1
            except Exception as e:
                logger.error(f"Error reading file {file_path}: {e}")
        
        logger.info(f"Processed {new_files_count} new files.")

    def _read_file(self, file_path):
        """
        Read content from a file based on extension.
        """
        ext = os.path.splitext(file_path)[1].lower()
        try:
            if ext in ['.txt', '.md']:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    return f.read()
            elif ext == '.json':
                with open(file_path, 'r', encoding='utf-8') as f:
                    data = json.load(f)
                    return json.dumps(data) # Convert back to string for processing
            elif ext == '.csv':
                import pandas as pd
                df = pd.read_csv(file_path)
                return df.to_string()
            elif ext == '.pdf':
                from pypdf import PdfReader
                reader = PdfReader(file_path)
                text = ""
                for page in reader.pages:
                    text += page.extract_text() + "\n"
                return text
            elif ext in ['.xlsx', '.xls']:
                import pandas as pd
                try:
                    # Read all sheets
                    dfs = pd.read_excel(file_path, sheet_name=None)
                    text = ""
                    for sheet_name, df in dfs.items():
                        text += f"Sheet: {sheet_name}\n{df.to_string()}\n"
                    return text
                except Exception as e:
                    logger.warning(f"Pandas failed to read {file_path}: {e}. Attempting fallback...")
                    # Fallback: Extract shared strings via zipfile (robust to style errors)
                    try:
                        import zipfile
                        import xml.etree.ElementTree as ET
                        with zipfile.ZipFile(file_path, 'r') as z:
                            if 'xl/sharedStrings.xml' in z.namelist():
                                with z.open('xl/sharedStrings.xml') as f:
                                    tree = ET.parse(f)
                                    root = tree.getroot()
                                    texts = []
                                    for elem in root.iter():
                                        if elem.tag.endswith('}t'):
                                            if elem.text:
                                                texts.append(elem.text)
                                    return "\n".join(texts)
                    except Exception as zip_e:
                        logger.error(f"Fallback failed for {file_path}: {zip_e}")
                        return None
            elif ext == '.pptx':
                from pptx import Presentation
                prs = Presentation(file_path)
                text = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
                return text
            elif ext == '.docx':
                from docx import Document
                doc = Document(file_path)
                text = ""
                for para in doc.paragraphs:
                    text += para.text + "\n"
                return text
            elif ext == '.zip':
                import zipfile
                import tempfile
                import shutil
                
                text = ""
                with zipfile.ZipFile(file_path, 'r') as zip_ref:
                    # Extract to a temp dir
                    with tempfile.TemporaryDirectory() as temp_dir:
                        zip_ref.extractall(temp_dir)
                        # Recursively read files in temp dir
                        for root, dirs, files in os.walk(temp_dir):
                            for file in files:
                                inner_path = os.path.join(root, file)
                                # Avoid infinite recursion if zip contains zip (simplified: just read content)
                                # We reuse _read_file but need to be careful about path
                                inner_content = self._read_file(inner_path)
                                if inner_content:
                                    text += f"\n--- File inside zip: {file} ---\n"
                                    text += inner_content
                return text
        except Exception as e:
            logger.error(f"Failed to parse {file_path}: {e}")
            return None
        
        return None

    def process_content(self, content, source):
        """
        Process a single document content: extract info, build graph, vectorize.
        """
        logger.info(f"Processing document: {source}")
        
        # 1. Vectorize
        self.vector_store.add_documents([content], metadata=[source])
        
        # 2. Extract Information
        triples = self.extractor.extract_triples(content)
        
        # 3. Build Graph
        self.graph_builder.add_triples(triples)

    def query(self, query_text):
        """
        Query the system.
        """
        results = {}
        
        # 1. Vector Search
        vector_results = self.vector_store.search(query_text)
        results['vector_search'] = vector_results
        
        # 2. Graph/Reasoning (Simplified: extract entities from query and look them up)
        query_entities = self.extractor.extract_entities(query_text)
        graph_results = []
        for entity, label in query_entities:
            # Look for connections (both ways)
            neighbors = self.graph_builder.query_graph(entity, direction='both')
            for neighbor, rel, direction in neighbors:
                if direction == 'out':
                    graph_results.append((entity, rel, neighbor))
                else:
                    graph_results.append((neighbor, rel, entity))
            
            # Look for transitive connections (example)
            # inferred = self.reasoner.infer_transitive(entity, "related_to")
            # if inferred:
            #     graph_results.append({"inferred_from": entity, "nodes": inferred})
                
        results['graph_search'] = graph_results
        
        return results

    def save_data(self):
        """
        Save all system data.
        """
        self.graph_builder.save_graph(os.path.join(self.storage_dir, "graph.pkl"))
        self.vector_store.save_index(os.path.join(self.storage_dir, "vector_index"))
        
        # Save processed files list
        with open(os.path.join(self.storage_dir, "processed_files.json"), 'w', encoding='utf-8') as f:
            json.dump(list(self.processed_files), f)
            
        logger.info("System data saved.")

    def load_data(self):
        """
        Load all system data.
        """
        self.graph_builder.load_graph(os.path.join(self.storage_dir, "graph.pkl"))
        self.vector_store.load_index(os.path.join(self.storage_dir, "vector_index"))
        
        # Load processed files list
        processed_path = os.path.join(self.storage_dir, "processed_files.json")
        if os.path.exists(processed_path):
            with open(processed_path, 'r', encoding='utf-8') as f:
                self.processed_files = set(json.load(f))
            logger.info(f"Loaded {len(self.processed_files)} processed files.")

if __name__ == "__main__":
    # Test
    system = KnowledgeSystem(storage_dir="d:/note/code/py/agent/data")
    # system.load_documents("d:/note/code/py/document") # Uncomment to run on real data
    
    # Dummy test
    system.process_content("Google is a tech giant. Google owns YouTube.", "dummy_source")
    system.save_data()
    print(system.query("What does Google own?"))
